"""Presentation extractor using python-pptx plus PaddleOCR for embedded images.

This is a lightweight structural extractor, similar to the DOCX path. It does
not render slides, so users should prefer manually exporting presentations to
PDF when layout, formulas, diagrams, or visual fidelity matter.
"""

from __future__ import annotations

import logging
import re
import tempfile
import zipfile
from pathlib import Path

from ingest.ocr_engine import get_engine

logger = logging.getLogger(__name__)


def extract(presentation_path: str | Path) -> str:
    """Extract text from PPTX, with a best-effort fallback for legacy PPT."""
    presentation_path = Path(presentation_path)
    if not presentation_path.exists():
        raise FileNotFoundError(f"Presentation not found: {presentation_path}")

    suffix = presentation_path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(presentation_path)
    if suffix == ".ppt":
        logger.warning("Legacy PPT extraction is best-effort only: %s", presentation_path.name)
        return _extract_legacy_ppt_text(presentation_path)
    raise ValueError(f"Unsupported presentation format: {presentation_path.suffix}")


def _extract_pptx(pptx_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "PPTX extraction requires python-pptx. Reinstall tree with current dependencies."
        ) from exc

    deck = Presentation(str(pptx_path))
    parts: list[str] = []

    for slide_index, slide in enumerate(deck.slides, start=1):
        slide_parts = _extract_slide_text(slide)
        if slide_parts:
            parts.append(f"# Slide {slide_index}\n\n" + "\n\n".join(slide_parts))

    image_texts = _extract_pptx_images(pptx_path)
    if image_texts:
        parts.append("# OCR from embedded images\n\n" + "\n\n".join(image_texts))

    return "\n\n".join(parts)


def _extract_slide_text(slide) -> list[str]:
    parts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                parts.append(text)
        if getattr(shape, "has_table", False):
            table_text = _table_to_markdown(shape.table)
            if table_text:
                parts.append(table_text)

    if getattr(slide, "has_notes_slide", False):
        notes_frame = slide.notes_slide.notes_text_frame
        notes = notes_frame.text.strip() if notes_frame else ""
        if notes:
            parts.append(f"Notes:\n{notes}")
    return parts


def _table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    if not rows:
        return ""
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
    body = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(part for part in (header, sep, body) if part)


def _extract_pptx_images(pptx_path: Path) -> list[str]:
    results: list[str] = []
    engine = get_engine()
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            media_names = sorted(
                name
                for name in archive.namelist()
                if name.startswith("ppt/media/") and Path(name).suffix.lower()
            )
            for index, name in enumerate(media_names, start=1):
                suffix = Path(name).suffix or ".png"
                tmp_path: Path | None = None
                try:
                    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                        tmp.write(archive.read(name))
                        tmp_path = Path(tmp.name)
                    logger.info("OCR-ing embedded presentation image #%d from %s", index, pptx_path.name)
                    text = engine.ocr_file(tmp_path)
                    if text.strip():
                        results.append(text)
                except Exception:
                    logger.exception("Failed to OCR embedded presentation image #%d", index)
                finally:
                    if tmp_path is not None:
                        try:
                            tmp_path.unlink()
                        except Exception:
                            pass
    except zipfile.BadZipFile:
        logger.warning("Invalid PPTX zip container: %s", pptx_path)
    return results


def _extract_legacy_ppt_text(ppt_path: Path) -> str:
    data = ppt_path.read_bytes()
    candidates = _printable_sequences(data.decode("utf-16le", errors="ignore"))
    candidates.extend(_printable_sequences(data.decode("latin-1", errors="ignore")))

    seen = set()
    lines = []
    for candidate in candidates:
        normalized = re.sub(r"\s+", " ", candidate).strip()
        if normalized and normalized not in seen:
            seen.add(normalized)
            lines.append(normalized)

    if not lines:
        logger.warning("No text extracted from legacy PPT: %s", ppt_path.name)
        return ""
    return "\n".join(lines)


def _printable_sequences(text: str) -> list[str]:
    return [
        match.group(0).strip()
        for match in re.finditer(r"[\w\u4e00-\u9fff，。；：！？、（）《》【】\s.,;:!?()\-+/=]{8,}", text)
        if match.group(0).strip()
    ]
